from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
MUTED_GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_NAME = "Calibri"
TITLE_SIZE = Pt(28)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(11)
HEADING_SIZE = Pt(14)

LEFT_MARGIN = Inches(0.8)
TOP_MARGIN = Inches(1.5)
CONTENT_WIDTH = Inches(11.7)
CONTENT_HEIGHT = Inches(5.5)

FLAG_LABELS = {
    "EXPLICIT": "Explicit (100%)",
    "PATTERN": "Pattern (75%)",
    "INFERRED": "Inferred (40%)",
}


@dataclass
class RenderContext:
    deck_name: str
    data_source_filename: str
    narrative_text: str
    narrative_confidence: float
    story_angle: str
    viz_recommendation: dict | None = None
    assumptions: list[dict] = field(default_factory=list)
    questions_and_answers: list[dict] = field(default_factory=list)
    quality_notes: list[dict] = field(default_factory=list)
    reconciliation_summary: dict = field(default_factory=dict)
    verified_at: str = ""


class VerificationGateError(Exception):
    """Raised when reconciliation has unresolved failures."""


def build_pptx(context: RenderContext) -> bytes:
    summary = context.reconciliation_summary
    failed = summary.get("failed_count", 0)
    dismissed = summary.get("dismissed_count", 0)
    if failed > 0:
        raise VerificationGateError(
            f"Render blocked: {failed} reconciliation check(s) failed. "
            f"Resolve or dismiss all failures before rendering."
        )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, context)
    _add_executive_summary_slide(prs, context)
    _add_data_slides(prs, context)
    _add_assumptions_slide(prs, context)
    _add_qa_slide(prs, context)
    _add_appendix_slide(prs, context)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_slide_notes(slide, context: RenderContext) -> None:
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        f"confidence: {context.narrative_confidence:.2f}, "
        f"assumptions: {len(context.assumptions)}"
    )


_MD_BOLD = re.compile(r"\*\*(.+?)\*\*")


def _strip_markdown(text: str) -> str:
    return _MD_BOLD.sub(r"\1", text)


def _set_text_props(run, size=BODY_SIZE, color=DARK_TEXT, bold=False):
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    run.font.bold = bold


def _add_title_slide(prs: Presentation, context: RenderContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = context.deck_name
    _set_text_props(run, size=Pt(36), color=DARK_TEXT, bold=True)

    sub_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(3.8), CONTENT_WIDTH, Inches(0.6))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"{datetime.now(timezone.utc).strftime('%B %d, %Y')}  |  Source: {context.data_source_filename}"
    _set_text_props(run2, size=SUBTITLE_SIZE, color=MUTED_GRAY)

    _add_slide_notes(slide, context)


def _add_executive_summary_slide(prs: Presentation, context: RenderContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
    tf_h = header_box.text_frame
    p_h = tf_h.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = "Executive Summary"
    _set_text_props(run_h, size=TITLE_SIZE, color=DARK_TEXT, bold=True)

    angle_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(1.3), CONTENT_WIDTH, Inches(0.5))
    tf_a = angle_box.text_frame
    p_a = tf_a.paragraphs[0]
    run_a = p_a.add_run()
    run_a.text = f"Story Angle: {(context.story_angle or '').replace('_', ' ').title()}"
    _set_text_props(run_a, size=SUBTITLE_SIZE, color=ACCENT_BLUE, bold=True)

    clean_text = _strip_markdown((context.narrative_text or "").replace("\n", " "))
    sentences = clean_text.split(". ")
    summary = ". ".join(sentences[:2]).strip()
    if summary and not summary.endswith("."):
        summary += "."

    body_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(2.0), CONTENT_WIDTH, CONTENT_HEIGHT)
    tf_b = body_box.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    run_b = p_b.add_run()
    run_b.text = summary
    _set_text_props(run_b, size=HEADING_SIZE, color=DARK_TEXT)

    conf_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.0), CONTENT_WIDTH, Inches(0.4))
    tf_c = conf_box.text_frame
    p_c = tf_c.paragraphs[0]
    run_c = p_c.add_run()
    run_c.text = f"Confidence: {context.narrative_confidence:.0%}"
    _set_text_props(run_c, size=BODY_SIZE, color=MUTED_GRAY)

    _add_slide_notes(slide, context)


def _add_data_slides(prs: Presentation, context: RenderContext) -> None:
    viz = context.viz_recommendation
    if not viz:
        return

    items = viz if isinstance(viz, list) else [viz]
    footnotes = [
        n.get("description", "") for n in context.quality_notes
        if n.get("severity") in ("warning", "info")
    ]
    footnote_text = " | ".join(footnotes[:3]) if footnotes else ""

    for item in items:
        chart_type = item.get("chart_type", "Chart") if isinstance(item, dict) else str(item)
        justification = item.get("justification", "") if isinstance(item, dict) else ""

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
        tf_h = header_box.text_frame
        p_h = tf_h.paragraphs[0]
        run_h = p_h.add_run()
        run_h.text = f"Data Visualization: {chart_type}"
        _set_text_props(run_h, size=TITLE_SIZE, color=DARK_TEXT, bold=True)

        placeholder_box = slide.shapes.add_textbox(
            Inches(2.0), Inches(2.5), Inches(9.3), Inches(3.0)
        )
        tf_p = placeholder_box.text_frame
        tf_p.word_wrap = True
        p_p = tf_p.paragraphs[0]
        p_p.alignment = PP_ALIGN.CENTER
        run_p = p_p.add_run()
        run_p.text = f"[Chart: {chart_type}]"
        _set_text_props(run_p, size=Pt(24), color=MUTED_GRAY)

        if justification:
            p_j = tf_p.add_paragraph()
            p_j.alignment = PP_ALIGN.CENTER
            run_j = p_j.add_run()
            run_j.text = justification
            _set_text_props(run_j, size=BODY_SIZE, color=MUTED_GRAY)

        if footnote_text:
            fn_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(6.8), CONTENT_WIDTH, Inches(0.4))
            tf_fn = fn_box.text_frame
            p_fn = tf_fn.paragraphs[0]
            run_fn = p_fn.add_run()
            run_fn.text = f"Note: {footnote_text}"
            _set_text_props(run_fn, size=Pt(9), color=MUTED_GRAY)

        _add_slide_notes(slide, context)


def _add_assumptions_slide(prs: Presentation, context: RenderContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
    tf_h = header_box.text_frame
    p_h = tf_h.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = "Assumptions & Inference Flags"
    _set_text_props(run_h, size=TITLE_SIZE, color=DARK_TEXT, bold=True)

    body_box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf_b = body_box.text_frame
    tf_b.word_wrap = True

    grouped: dict[str, list[dict]] = {"EXPLICIT": [], "PATTERN": [], "INFERRED": []}
    for a in context.assumptions:
        flag = a.get("flag_type", "EXPLICIT")
        grouped.setdefault(flag, []).append(a)

    first = True
    for flag_type in ("EXPLICIT", "PATTERN", "INFERRED"):
        items = grouped.get(flag_type, [])
        if not items:
            continue

        if first:
            p = tf_b.paragraphs[0]
            first = False
        else:
            p = tf_b.add_paragraph()
            p.space_before = Pt(12)

        run = p.add_run()
        run.text = FLAG_LABELS.get(flag_type, flag_type)
        _set_text_props(run, size=HEADING_SIZE, color=ACCENT_BLUE, bold=True)

        for a in items:
            p_item = tf_b.add_paragraph()
            p_item.level = 1
            run_item = p_item.add_run()
            conf = a.get("confidence", 0)
            run_item.text = f"• {a.get('text', '')}  [{conf:.0%}]"
            _set_text_props(run_item, size=BODY_SIZE, color=DARK_TEXT)

    if first:
        p = tf_b.paragraphs[0]
        run = p.add_run()
        run.text = "No assumptions flagged."
        _set_text_props(run, size=BODY_SIZE, color=MUTED_GRAY)

    _add_slide_notes(slide, context)


def _add_qa_slide(prs: Presentation, context: RenderContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
    tf_h = header_box.text_frame
    p_h = tf_h.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = "Questions & Answers"
    _set_text_props(run_h, size=TITLE_SIZE, color=DARK_TEXT, bold=True)

    body_box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf_b = body_box.text_frame
    tf_b.word_wrap = True

    qa_pairs = context.questions_and_answers
    if not qa_pairs:
        p = tf_b.paragraphs[0]
        run = p.add_run()
        run.text = "No Q&A data available."
        _set_text_props(run, size=BODY_SIZE, color=MUTED_GRAY)
        _add_slide_notes(slide, context)
        return

    first = True
    for qa in qa_pairs:
        question = qa.get("question", qa.get("template", ""))
        answer = qa.get("answer", qa.get("raw_answer", "N/A"))

        if first:
            p_q = tf_b.paragraphs[0]
            first = False
        else:
            p_q = tf_b.add_paragraph()
            p_q.space_before = Pt(10)

        run_q = p_q.add_run()
        run_q.text = f"Q: {question}"
        _set_text_props(run_q, size=BODY_SIZE, color=ACCENT_BLUE, bold=True)

        p_a = tf_b.add_paragraph()
        run_a = p_a.add_run()
        run_a.text = f"A: {answer}"
        _set_text_props(run_a, size=BODY_SIZE, color=DARK_TEXT)

    _add_slide_notes(slide, context)


def _add_appendix_slide(prs: Presentation, context: RenderContext) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
    tf_h = header_box.text_frame
    p_h = tf_h.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = "Appendix"
    _set_text_props(run_h, size=TITLE_SIZE, color=DARK_TEXT, bold=True)

    body_box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf_b = body_box.text_frame
    tf_b.word_wrap = True

    # Data Quality Notes
    p_dq = tf_b.paragraphs[0]
    run_dq = p_dq.add_run()
    run_dq.text = "Data Quality Notes"
    _set_text_props(run_dq, size=HEADING_SIZE, color=ACCENT_BLUE, bold=True)

    if context.quality_notes:
        for note in context.quality_notes[:5]:
            p_n = tf_b.add_paragraph()
            p_n.level = 1
            run_n = p_n.add_run()
            severity = note.get("severity", "info")
            desc = note.get("description", str(note))
            run_n.text = f"• [{severity.upper()}] {desc}"
            _set_text_props(run_n, size=BODY_SIZE, color=DARK_TEXT)
    else:
        p_n = tf_b.add_paragraph()
        run_n = p_n.add_run()
        run_n.text = "No quality issues detected."
        _set_text_props(run_n, size=BODY_SIZE, color=MUTED_GRAY)

    # Reconciliation Summary
    p_rs = tf_b.add_paragraph()
    p_rs.space_before = Pt(16)
    run_rs = p_rs.add_run()
    run_rs.text = "Reconciliation Summary"
    _set_text_props(run_rs, size=HEADING_SIZE, color=ACCENT_BLUE, bold=True)

    summary = context.reconciliation_summary
    passed = summary.get("passed_count", 0)
    failed = summary.get("failed_count", 0)
    dismissed = summary.get("dismissed_count", 0)
    total = summary.get("total_checks", passed + failed + dismissed)

    p_sum = tf_b.add_paragraph()
    p_sum.level = 1
    run_sum = p_sum.add_run()
    run_sum.text = f"• {passed}/{total} checks passed"
    _set_text_props(run_sum, size=BODY_SIZE, color=DARK_TEXT)

    if dismissed > 0:
        p_dis = tf_b.add_paragraph()
        p_dis.level = 1
        run_dis = p_dis.add_run()
        run_dis.text = f"• {dismissed} check(s) dismissed with reason"
        _set_text_props(run_dis, size=BODY_SIZE, color=DARK_TEXT)

    # Verification Timestamp
    p_vt = tf_b.add_paragraph()
    p_vt.space_before = Pt(16)
    run_vt = p_vt.add_run()
    run_vt.text = "Verification Timestamp"
    _set_text_props(run_vt, size=HEADING_SIZE, color=ACCENT_BLUE, bold=True)

    p_ts = tf_b.add_paragraph()
    p_ts.level = 1
    run_ts = p_ts.add_run()
    run_ts.text = f"• Verified at: {context.verified_at or 'N/A'}"
    _set_text_props(run_ts, size=BODY_SIZE, color=DARK_TEXT)

    _add_slide_notes(slide, context)
