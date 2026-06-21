import io

from pptx import Presentation

from app.services.render.pptx_builder import build_pptx, RenderContext


def _make_context(**overrides) -> RenderContext:
    defaults = {
        "deck_name": "Q3 Financial Review",
        "data_source_filename": "financials.csv",
        "narrative_text": "Revenue grew 15% YoY. Cost increased 8% driven by hiring. Margins improved overall.",
        "narrative_confidence": 0.85,
        "story_angle": "trend",
        "viz_recommendation": {"chart_type": "bar", "justification": "Shows trend over time"},
        "assumptions": [
            {"text": "Revenue is $5M", "flag_type": "EXPLICIT", "confidence": 1.0, "source_reference": "data"},
            {"text": "Growth trend detected", "flag_type": "PATTERN", "confidence": 0.75, "source_reference": "analysis"},
            {"text": "Market stable", "flag_type": "INFERRED", "confidence": 0.4, "source_reference": "context"},
        ],
        "questions_and_answers": [
            {"question": "What is the primary metric?", "answer": "Revenue growth"},
            {"question": "Who is the audience?", "answer": "Board of directors"},
        ],
        "quality_notes": [
            {"severity": "warning", "description": "5 duplicate rows removed"},
        ],
        "reconciliation_summary": {
            "total_checks": 5,
            "passed_count": 4,
            "failed_count": 0,
            "dismissed_count": 1,
        },
        "verified_at": "2026-06-21T14:30:00",
    }
    defaults.update(overrides)
    return RenderContext(**defaults)


def test_build_pptx_returns_valid_bytes():
    ctx = _make_context()
    result = build_pptx(ctx)
    assert isinstance(result, bytes)
    assert len(result) > 0
    assert result[:2] == b"PK"


def test_slide_count_with_viz():
    ctx = _make_context()
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    # title + summary + 1 data slide + assumptions + Q&A + appendix = 6
    assert len(prs.slides) == 6


def test_slide_count_without_viz():
    ctx = _make_context(viz_recommendation=None)
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    # title + summary + assumptions + Q&A + appendix = 5 (no data slides)
    assert len(prs.slides) == 5


def test_title_slide_contains_deck_name():
    ctx = _make_context(deck_name="Test Deck Alpha")
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert any("Test Deck Alpha" in t for t in texts)


def test_title_slide_contains_source_filename():
    ctx = _make_context(data_source_filename="my_data.csv")
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert any("my_data.csv" in t for t in texts)


def test_executive_summary_slide():
    ctx = _make_context(narrative_text="First sentence. Second sentence. Third sentence.")
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[1]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    full_text = " ".join(texts)
    assert "Executive Summary" in full_text
    assert "First sentence" in full_text


def test_assumptions_slide_groups_by_flag_type():
    ctx = _make_context()
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    # With viz: slides are title(0), summary(1), data(2), assumptions(3), Q&A(4), appendix(5)
    slide = prs.slides[3]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    full_text = " ".join(texts)
    assert "Assumptions & Inference Flags" in full_text
    assert "Explicit (100%)" in full_text
    assert "Pattern (75%)" in full_text
    assert "Inferred (40%)" in full_text


def test_qa_slide_contains_pairs():
    ctx = _make_context()
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[4]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    full_text = " ".join(texts)
    assert "Questions & Answers" in full_text
    assert "What is the primary metric?" in full_text
    assert "Revenue growth" in full_text


def test_appendix_slide_contains_quality_notes():
    ctx = _make_context()
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[5]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    full_text = " ".join(texts)
    assert "Appendix" in full_text
    assert "5 duplicate rows removed" in full_text
    assert "4/5 checks passed" in full_text
    assert "2026-06-21" in full_text


def test_slide_notes_contain_metadata():
    ctx = _make_context()
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    for slide in prs.slides:
        notes = slide.notes_slide.notes_text_frame.text
        assert "confidence: 0.85" in notes
        assert "assumptions: 3" in notes


def test_multiple_viz_recommendations():
    ctx = _make_context(viz_recommendation=[
        {"chart_type": "bar", "justification": "Comparison"},
        {"chart_type": "line", "justification": "Trend"},
    ])
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    # title + summary + 2 data slides + assumptions + Q&A + appendix = 7
    assert len(prs.slides) == 7


def test_empty_assumptions():
    ctx = _make_context(assumptions=[])
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_empty_qa():
    ctx = _make_context(questions_and_answers=[])
    result = build_pptx(ctx)
    prs = Presentation(io.BytesIO(result))
    # Q&A slide still exists but shows "No Q&A data available."
    slide = prs.slides[4] if ctx.viz_recommendation else prs.slides[3]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    full_text = " ".join(texts)
    assert "No Q&A data available" in full_text
